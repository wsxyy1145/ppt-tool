"""
PPT工具 - 添加图片和图表
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
import os


def add_image_to_slide(image_path: str, slide_index: int = 0, left: float = 1, top: float = 2, width: float = 4,
                       output_path: str = "output.pptx"):
    """
    在指定幻灯片上添加图片
    
    参数:
        image_path: 图片路径
        slide_index: 幻灯片索引（从0开始）
        left: 左边距（英寸）
        top: 上边距（英寸）
        width: 图片宽度（英寸）
        output_path: 输出文件路径
    """
    prs = Presentation()
    
    # 确保有足够的幻灯片
    while len(prs.slides) <= slide_index:
        prs.slides.add_slide(prs.slide_layouts[1])
    
    slide = prs.slides[slide_index]
    
    # 添加图片
    pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top), width=Inches(width))
    
    prs.save(output_path)
    return os.path.abspath(output_path)


def add_chart_to_slide(chart_data: dict, chart_type: str = "column", slide_index: int = 0,
                       left: float = 1, top: float = 2, width: float = 6, height: float = 4.5,
                       title: str = "图表", output_path: str = "output.pptx"):
    """
    在指定幻灯片上添加图表
    
    参数:
        chart_data: 图表数据，格式如：
            {
                "categories": ["第一季度", "第二季度", "第三季度", "第四季度"],
                "series": {"销售额": (100, 120, 140, 180), "利润": (30, 40, 50, 60)}
            }
        chart_type: 图表类型 column, line, pie, bar
        slide_index: 幻灯片索引
        left, top, width, height: 位置和大小
        title: 图表标题
        output_path: 输出文件路径
    """
    # 映射图表类型
    chart_type_map = {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED
    }
    xl_chart_type = chart_type_map.get(chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
    
    prs = Presentation()
    
    # 确保有足够的幻灯片
    while len(prs.slides) <= slide_index:
        prs.slides.add_slide(prs.slide_layouts[1])
    
    slide = prs.slides[slide_index]
    
    # 设置标题
    slide.shapes.title.text = title
    
    # 创建图表数据
    chart_data_obj = CategoryChartData()
    chart_data_obj.categories = chart_data["categories"]
    
    for series_name, values in chart_data["series"].items():
        chart_data_obj.add_series(series_name, values)
    
    # 添加图表
    x, y, cx, cy = Inches(left), Inches(top), Inches(width), Inches(height)
    chart = slide.shapes.add_chart(xl_chart_type, x, y, cx, cy, chart_data_obj).chart
    
    prs.save(output_path)
    return os.path.abspath(output_path)


def add_shape(shape_type: str = "rectangle", slide_index: int = 0, left: float = 1, top: float = 1,
              width: float = 2, height: float = 1, text: str = "", output_path: str = "output.pptx"):
    """
    在指定幻灯片上添加形状
    
    参数:
        shape_type: 形状类型 rectangle, oval, triangle, etc.
        slide_index: 幻灯片索引
        left, top, width, height: 位置和大小
        text: 形状内的文字
        output_path: 输出文件路径
    """
    shape_type_map = {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "oval": MSO_SHAPE.OVAL,
        "triangle": MSO_SHAPE.TRIANGLE,
        "arrow": MSO_SHAPE.ARROW,
        "star": MSO_SHAPE.STAR_5
    }
    ms_shape = shape_type_map.get(shape_type.lower(), MSO_SHAPE.RECTANGLE)
    
    prs = Presentation()
    
    while len(prs.slides) <= slide_index:
        prs.slides.add_slide(prs.slide_layouts[1])
    
    slide = prs.slides[slide_index]
    
    shape = slide.shapes.add_shape(ms_shape, Inches(left), Inches(top), Inches(width), Inches(height))
    
    if text:
        shape.text = text
    
    prs.save(output_path)
    return os.path.abspath(output_path)


if __name__ == '__main__':
    print("创建带图表的PPT...")
    
    # 示例：创建带图表的PPT
    chart_data = {
        "categories": ["第一季度", "第二季度", "第三季度", "第四季度"],
        "series": {
            "销售额": (100, 120, 140, 180),
            "利润": (30, 40, 50, 60)
        }
    }
    
    add_chart_to_slide(
        chart_data=chart_data,
        chart_type="column",
        title="年度销售数据",
        output_path="chart_example.pptx"
    )
    
    print("完成！")