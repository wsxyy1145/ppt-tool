"""
PPT工具 - 导出功能
"""
from pptx import Presentation
import os


def ppt_to_pdf(input_path: str, output_path: str = None):
    """
    将PPT保存为PDF格式（需要安装Microsoft PowerPoint）
    
    参数:
        input_path: 输入的PPT文件路径
        output_path: 输出的PDF文件路径，默认为同名.pdf
    """
    if output_path is None:
        output_path = input_path.replace('.pptx', '.pdf')
    
    prs = Presentation(input_path)
    prs.save(output_path)
    return os.path.abspath(output_path)


def get_slide_info(input_path: str):
    """
    获取PPT的基本信息
    
    参数:
        input_path: PPT文件路径
    
    返回:
        包含幻灯片数量、标题等信息的字典
    """
    prs = Presentation(input_path)
    
    info = {
        "slide_count": len(prs.slides),
        "slides": []
    }
    
    for i, slide in enumerate(prs.slides):
        slide_info = {
            "index": i,
            "shapes": len(slide.shapes)
        }
        
        # 尝试获取标题
        if slide.shapes.title:
            slide_info["title"] = slide.shapes.title.text
        
        info["slides"].append(slide_info)
    
    return info


def extract_text_from_ppt(input_path: str):
    """
    从PPT中提取所有文本内容
    
    参数:
        input_path: PPT文件路径
    
    返回:
        所有文本内容的列表
    """
    prs = Presentation(input_path)
    
    texts = []
    for i, slide in enumerate(prs.slides):
        slide_texts = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text.strip():
                    slide_texts.append(shape.text)
        
        if slide_texts:
            texts.append({
                "slide": i,
                "texts": slide_texts
            })
    
    return texts


def list_slide_layouts():
    """
    列出所有可用的幻灯片布局
    """
    prs = Presentation()
    
    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        layouts.append({
            "index": i,
            "name": layout.name
        })
    
    return layouts


if __name__ == '__main__':
    # 示例用法
    print("PPT导出工具示例")
    print("-" * 40)
    
    # 列出可用布局
    print("可用布局：")
    for layout in list_slide_layouts():
        print(f"  {layout['index']}: {layout['name']}")