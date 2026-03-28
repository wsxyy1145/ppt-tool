"""
PPT工具 - 快速创建演示文稿
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os


def create_simple_ppt(title: str, subtitle: str = "", slides_content: list = None, output_path: str = "output.pptx"):
    """
    快速创建一个简单的PPT
    
    参数:
        title: 演示文稿标题
        subtitle: 副标题
        slides_content: 幻灯片内容列表，每个元素可以是字符串或字典
        output_path: 输出文件路径
    
    示例:
        create_simple_ppt(
            title="年度报告",
            subtitle="2024年度总结",
            slides_content=[
                "第一页内容",
                {"title": "第二页标题", "content": "第二页内容"},
                "第三页内容"
            ]
        )
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 封面页
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    
    # 添加内容页
    if slides_content:
        for content in slides_content:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            
            if isinstance(content, str):
                # 纯文本内容
                slide.shapes.title.text = "内容"
                slide.placeholders[1].text_frame.text = content
            elif isinstance(content, dict):
                # 字典格式：包含title和content
                if "title" in content:
                    slide.shapes.title.text = content["title"]
                if "content" in content:
                    slide.placeholders[1].text_frame.text = content["content"]
    
    # 保存
    prs.save(output_path)
    return os.path.abspath(output_path)


def create_title_slide(title: str, subtitle: str = "", output_path: str = "title_slide.pptx"):
    """创建标题页"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    prs.save(output_path)
    return os.path.abspath(output_path)


def create_content_slide(title: str, content: str, output_path: str = "content_slide.pptx"):
    """创建标题+内容页"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text_frame.text = content
    prs.save(output_path)
    return os.path.abspath(output_path)


if __name__ == '__main__':
    # 示例用法
    print("创建示例PPT...")
    
    # 示例1：快速创建
    create_simple_ppt(
        title="产品发布会",
        subtitle="2024年度新品",
        slides_content=[
            "这是第一页的内容",
            {"title": "产品特点", "content": "• 特点1\n• 特点2\n• 特点3"},
            {"title": "市场价格", "content": "¥2999起"}
        ]
    )
    
    print("完成！")